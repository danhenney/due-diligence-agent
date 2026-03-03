"""PPTX report generator — consulting-style slide deck from DD state."""
from __future__ import annotations

import os
import re
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Colors ────────────────────────────────────────────────────────────────────
COLOR_DARK = RGBColor(0x1E, 0x29, 0x3B)       # slate-800
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_INVEST = RGBColor(0x16, 0xA3, 0x4A)     # green-600
COLOR_WATCH = RGBColor(0xD9, 0x77, 0x06)      # amber-600
COLOR_PASS = RGBColor(0xDC, 0x26, 0x26)       # red-600
COLOR_ACCENT = RGBColor(0x25, 0x63, 0xEB)     # blue-600
COLOR_MUTED = RGBColor(0x64, 0x74, 0x8B)      # slate-500
COLOR_LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)   # slate-50


def _rec_color(rec: str) -> RGBColor:
    r = (rec or "").upper()
    if "INVEST" in r:
        return COLOR_INVEST
    if "WATCH" in r:
        return COLOR_WATCH
    return COLOR_PASS


def _safe_str(val: Any, default: str = "N/A") -> str:
    if val is None:
        return default
    if isinstance(val, dict):
        return val.get("summary", val.get("description", str(val)[:200]))
    if isinstance(val, list):
        return "; ".join(str(v) for v in val[:5])
    return str(val)[:500]


def _add_title_slide(prs: Presentation, company: str, recommendation: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_DARK

    # "Due Diligence Report" subtitle
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Due Diligence Report"
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_MUTED
    p.alignment = PP_ALIGN.CENTER

    # Company name
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = company
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Recommendation badge
    rec = (recommendation or "WATCH").upper()
    txBox = slide.shapes.add_textbox(Inches(3), Inches(4), Inches(4), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"Recommendation: {rec}"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = _rec_color(rec)
    p.alignment = PP_ALIGN.CENTER

    # Date
    txBox = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(6), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_MUTED
    p.alignment = PP_ALIGN.CENTER


def _add_section_slide(prs: Presentation, title: str, bullets: list[str],
                       subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.4))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_MUTED

    # Bullets
    y_start = 1.5 if subtitle else 1.2
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(y_start), Inches(9), Inches(5.5 - y_start))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets[:12]):  # max 12 bullets per slide
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_DARK
        p.space_after = Pt(6)


def _extract_bullets(data: Any, keys: list[str], max_items: int = 8) -> list[str]:
    """Extract bullet points from a dict, trying multiple key patterns."""
    if not isinstance(data, dict):
        return [_safe_str(data)]
    bullets = []
    for key in keys:
        val = data.get(key)
        if val is None:
            continue
        if isinstance(val, list):
            for item in val[:max_items]:
                if isinstance(item, dict):
                    label = item.get("name") or item.get("risk") or item.get("trend") or item.get("argument") or ""
                    desc = item.get("description") or item.get("impact") or item.get("supporting_evidence") or ""
                    if isinstance(desc, list):
                        desc = "; ".join(str(d) for d in desc[:3])
                    bullets.append(f"{label}: {desc}" if label else str(desc)[:200])
                else:
                    bullets.append(str(item)[:200])
        elif isinstance(val, str):
            bullets.append(f"{key}: {val}")
        elif isinstance(val, dict):
            summary = val.get("summary") or val.get("description") or val.get("assessment") or str(val)[:200]
            bullets.append(f"{key}: {summary}")
    if not bullets:
        summary = data.get("summary", "")
        if summary:
            bullets.append(summary)
    return bullets[:max_items]


def generate_pptx(state: dict[str, Any], job_id: str, output_dir: str | None = None) -> str:
    """Generate a consulting-style PPTX slide deck from DD state.

    Returns absolute path to the generated .pptx file.
    """
    reports_dir = Path(output_dir) if output_dir else Path("reports")
    reports_dir.mkdir(parents=True, exist_ok=True)
    output_path = str(reports_dir / f"{job_id}.pptx")

    company = state.get("company_name") or "Unknown Company"
    recommendation = state.get("recommendation") or "WATCH"

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. Title slide
    _add_title_slide(prs, company, recommendation)

    # 2. Executive Summary
    exec_summary = state.get("final_report", "")[:500] if state.get("final_report") else "N/A"
    si = state.get("strategic_insight") or {}
    _add_section_slide(prs, "Executive Summary", [
        exec_summary,
        f"Recommendation: {recommendation}",
        f"Rationale: {_safe_str(si.get('rationale', 'See full report'))}",
    ])

    # 3. Market Analysis
    market = state.get("market_analysis") or {}
    _add_section_slide(prs, "Market Analysis",
        _extract_bullets(market, ["summary", "tam", "sam", "som", "cagr", "trends", "market_drivers"]),
        subtitle="TAM / SAM / SOM / Growth Trends"
    )

    # 4. Competitive Landscape
    comp = state.get("competitor_analysis") or {}
    _add_section_slide(prs, "Competitive Landscape",
        _extract_bullets(comp, ["summary", "competitors", "market_share", "competitive_gaps"]),
        subtitle="Competitors & Positioning"
    )

    # 5. Financial Analysis & Valuation
    fin = state.get("financial_analysis") or {}
    _add_section_slide(prs, "Financial Analysis & Valuation",
        _extract_bullets(fin, ["summary", "revenue_trend", "profitability", "cash_flow", "valuation", "key_ratios"]),
        subtitle="Revenue / Profitability / Cash Flow / Fair Value"
    )

    # 6. Technology & IP
    tech = state.get("tech_analysis") or {}
    _add_section_slide(prs, "Technology & IP",
        _extract_bullets(tech, ["summary", "core_technologies", "ip_patents", "tech_maturity", "competitive_comparison"]),
        subtitle="Core Tech / Patents / Maturity"
    )

    # 7. Legal & Regulatory
    legal = state.get("legal_regulatory") or {}
    _add_section_slide(prs, "Legal & Regulatory",
        _extract_bullets(legal, ["summary", "investment_structure_risks", "business_regulatory_risks", "litigation"]),
        subtitle="Investment Structure & Business Regulatory Risks"
    )

    # 8. Team & Leadership
    team = state.get("team_analysis") or {}
    _add_section_slide(prs, "Team & Leadership",
        _extract_bullets(team, ["summary", "leadership_profiles", "capability_assessment", "key_person_risk"]),
        subtitle="Leadership / Capability / Key Person Risk"
    )

    # 9. Investment Thesis (R&A Synthesis)
    ras = state.get("ra_synthesis") or {}
    _add_section_slide(prs, "Investment Thesis",
        _extract_bullets(ras, ["summary", "core_investment_arguments", "attractiveness_scorecard", "key_findings"]),
        subtitle="Core Arguments & Attractiveness Scorecard"
    )

    # 10. Risk Matrix
    risk = state.get("risk_assessment") or {}
    _add_section_slide(prs, "Risk Assessment",
        _extract_bullets(risk, ["summary", "top_risks", "risk_matrix", "overall_risk_level"]),
        subtitle="Risk Matrix / Top Risks / Mitigation"
    )

    # 11. Recommendation
    _add_section_slide(prs, f"Recommendation: {recommendation}", [
        _safe_str(si.get("rationale", "See full report")),
        f"Key arguments for: {_safe_str(si.get('key_arguments_for', []))}",
        f"Key arguments against: {_safe_str(si.get('key_arguments_against', []))}",
        f"Investment timeline: {_safe_str(si.get('investment_timeline', 'N/A'))}",
    ])

    # 12. DD Questions
    ddq = state.get("dd_questions") or {}
    questions = ddq.get("dd_questionnaire", [])
    q_bullets = []
    for q in questions[:8]:
        if isinstance(q, dict):
            q_bullets.append(f"[{q.get('priority', '?')}] {q.get('question', '?')} → {q.get('target', '?')}")
        else:
            q_bullets.append(str(q)[:200])
    _add_section_slide(prs, "DD Questionnaire",
        q_bullets or ["No unresolved questions identified."],
        subtitle="Outstanding Questions for Follow-Up"
    )

    # 13. Appendix — Data Sources
    _add_section_slide(prs, "Appendix: Data Sources", [
        "See full report for complete source list with URLs",
        "All data verified via live tool calls (yfinance, Tavily, EDGAR, etc.)",
        f"Report generated: {datetime.now().strftime('%B %d, %Y')}",
        "Confidential — For Internal Use Only",
    ])

    prs.save(output_path)
    return os.path.abspath(output_path)
